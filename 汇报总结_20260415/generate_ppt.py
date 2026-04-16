#!/usr/bin/env python3
"""
Academic-style PPT generator for the lock-clock project.

Design goal:
- emphasize controlled baseline-vs-static evidence
- separate artifact-backed results from preserved historical summaries
- present predictor accuracy with explicit scope and limitations
"""

from __future__ import annotations

from dataclasses import dataclass
import math

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


# =========================
# Theme
# =========================

BG = RGBColor(250, 251, 253)
TITLE = RGBColor(25, 38, 59)
TEXT = RGBColor(46, 58, 74)
MUTED = RGBColor(108, 117, 125)
ACCENT = RGBColor(52, 86, 139)
ACCENT_LIGHT = RGBColor(235, 241, 249)
SUCCESS = RGBColor(36, 121, 93)
SUCCESS_LIGHT = RGBColor(235, 245, 240)
WARNING = RGBColor(168, 92, 44)
WARNING_LIGHT = RGBColor(252, 243, 235)
DANGER = RGBColor(160, 55, 70)
DANGER_LIGHT = RGBColor(249, 238, 241)
BORDER = RGBColor(213, 220, 230)


@dataclass(frozen=True)
class StaticPoint:
    label: str
    time_s: float
    avg_power_w: float
    energy_j: float


@dataclass(frozen=True)
class ControlledCase:
    title: str
    topology: str
    evidence_label: str
    baseline: StaticPoint
    static_points: tuple[StaticPoint, ...]


CASE_A = ControlledCase(
    title="案例 A：V100 双机 `TP=1, PP=4, DP=4`",
    topology="Megatron-DeepSpeed baseline 与 static 对照；Zeus 历史摘要保存在 memory-bank",
    evidence_label="证据等级 B：memory-bank/observability preserved Zeus summaries",
    baseline=StaticPoint("baseline", 1316.7, 3170.6, 4174627.4),
    static_points=(
        StaticPoint("1252 MHz", 1300.1, 2367.9, 3078625.8),
        StaticPoint("1260 MHz", 1287.9, 2382.0, 3067834.2),
        StaticPoint("1267 MHz", 1286.3, 2407.7, 3097007.8),
    ),
)

CASE_B = ControlledCase(
    title="案例 B：V100 双机 `TP=2, PP=2, DP=4`",
    topology="Megatron-DeepSpeed baseline 与 static 对照；Zeus 历史摘要保存在 memory-bank",
    evidence_label="证据等级 B：baseline summary in memory-bank, with partial local static artifacts",
    baseline=StaticPoint("baseline", 1030.6, 3686.1, 3798903.8),
    static_points=(
        StaticPoint("1072 MHz", 1053.6, 2349.8, 2475841.1),
        StaticPoint("1080 MHz", 1046.1, 2362.5, 2471497.2),
        StaticPoint("1087 MHz", 1040.0, 2386.6, 2482053.4),
        StaticPoint("1125 MHz", 1006.1, 2477.3, 2492545.3),
    ),
)

IB_METRICS = {
    "pair": "fresh formal replay: `2x4 -> 2x8`, IB, transport-consistent",
    "time_mape": 11.48,
    "power_mape": 3.28,
    "energy_mape": 7.86,
    "alpha_dp": "2.220525e-11 s/byte",
    "points": (
        ("990 MHz", "401.19", "435.40", "8.53%", "1189.17", "1128.43", "5.11%"),
        ("1080 MHz", "372.24", "415.90", "11.73%", "1257.39", "1229.19", "2.24%"),
        ("1155 MHz", "351.53", "401.42", "14.19%", "1353.14", "1319.64", "2.48%"),
    ),
}

ETH_METRICS = {
    "pair": "formal replay: `1x4 -> 2x4`, Ethernet, TP=1/PP=2 path",
    "time_mape": 5.16,
    "power_mape": 12.38,
    "energy_mape": 10.42,
    "alpha_dp": "7.321889e-10 s/byte",
    "points": (
        ("1005 MHz", "217.96", "227.29", "4.28%", "16.79%"),
        ("1200 MHz", "202.78", "217.12", "7.07%", "12.95%"),
        ("1395 MHz", "217.56", "208.56", "4.14%", "7.40%"),
    ),
}


def set_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def pct_change(value: float, baseline: float) -> float:
    return ((value - baseline) / baseline) * 100.0


def fmt_pct(value: float) -> str:
    return f"{value:+.1f}%"


def fmt_energy_kj(value_j: float) -> str:
    return f"{value_j / 1000.0:.1f}"


def add_textbox(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = TEXT,
    align=PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color


def add_bullets(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    bullets: list[str],
    *,
    size: int = 18,
    color: RGBColor = TEXT,
    bullet_prefix: str = "• ",
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for index, bullet in enumerate(bullets):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.text = f"{bullet_prefix}{bullet}"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = color
        paragraph.space_after = Pt(8)


def add_header(slide, title: str, subtitle: str | None = None) -> None:
    add_textbox(slide, 0.7, 0.35, 14.8, 0.55, title, size=28, bold=True, color=TITLE)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.02), Inches(1.2), Inches(0.05))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()
    if subtitle:
        add_textbox(slide, 0.7, 1.1, 14.5, 0.35, subtitle, size=12, color=MUTED)


def add_footer(slide, text: str, page_no: int) -> None:
    add_textbox(slide, 0.7, 8.45, 13.8, 0.2, text, size=9, color=MUTED)
    add_textbox(slide, 14.8, 8.42, 0.7, 0.2, str(page_no), size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def add_note_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    bullets: list[str],
    *,
    title_size: int = 15,
    bullet_size: int = 13,
    fill_color: RGBColor = ACCENT_LIGHT,
    border_color: RGBColor = BORDER,
) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    add_textbox(slide, x + 0.2, y + 0.15, w - 0.4, 0.35, title, size=title_size, bold=True, color=TITLE)
    add_bullets(slide, x + 0.2, y + 0.55, w - 0.4, h - 0.75, bullets, size=bullet_size)


def add_metric_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    value: str,
    label: str,
    note: str,
    *,
    value_color: RGBColor = ACCENT,
) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = BORDER
    add_textbox(slide, x, y + 0.28, w, 0.55, value, size=26, bold=True, color=value_color, align=PP_ALIGN.CENTER)
    add_textbox(slide, x, y + 0.92, w, 0.3, label, size=13, bold=True, color=TITLE, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.15, y + 1.22, w - 0.3, 0.42, note, size=10, color=MUTED, align=PP_ALIGN.CENTER)


def build_table(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    rows: list[list[str]],
    col_widths: list[float],
    *,
    header_font_size: int = 11,
    body_font_size: int = 12,
) -> None:
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h)).table
    for idx, width in enumerate(col_widths):
        table.columns[idx].width = Inches(width)

    row_height = h / len(rows)
    for row_idx in range(len(rows)):
        table.rows[row_idx].height = Inches(row_height)

    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT_LIGHT if row_idx == 0 else RGBColor(255, 255, 255)
            if row_idx > 0 and col_idx >= 2 and value.startswith("-"):
                cell.fill.fore_color.rgb = SUCCESS_LIGHT
            elif row_idx > 0 and col_idx == 2 and value.startswith("+"):
                cell.fill.fore_color.rgb = WARNING_LIGHT
            elif row_idx > 0 and col_idx >= 4 and value.startswith("+"):
                cell.fill.fore_color.rgb = DANGER_LIGHT
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = value
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(header_font_size if row_idx == 0 else body_font_size)
            p.font.bold = row_idx == 0
            p.font.color.rgb = TITLE if row_idx == 0 else TEXT


def add_chart_panel(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    *,
    subtitle: str | None = None,
) -> tuple[float, float, float, float]:
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
    panel.line.color.rgb = BORDER
    add_textbox(slide, x + 0.18, y + 0.12, w - 0.36, 0.28, title, size=13, bold=True, color=TITLE)
    if subtitle:
        add_textbox(slide, x + 0.18, y + 0.38, w - 0.36, 0.22, subtitle, size=9, color=MUTED)
        chart_y = y + 0.66
    else:
        chart_y = y + 0.48
    return x + 0.12, chart_y, w - 0.24, h - (chart_y - y) - 0.12


def add_column_chart(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    categories: list[str],
    series_defs: list[tuple[str, list[float], RGBColor]],
    *,
    subtitle: str | None = None,
    value_axis_format: str = '0.0',
    y_min: float | None = None,
    y_max: float | None = None,
    show_legend: bool = True,
) -> None:
    chart_x, chart_y, chart_w, chart_h = add_chart_panel(slide, x, y, w, h, title, subtitle=subtitle)
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values, _color in series_defs:
        chart_data.add_series(name, values)

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(chart_x),
        Inches(chart_y),
        Inches(chart_w),
        Inches(chart_h),
        chart_data,
    )
    chart = graphic_frame.chart
    chart.has_legend = show_legend
    if show_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    plot = chart.plots[0]
    plot.gap_width = 70
    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(10)
    value_axis = chart.value_axis
    value_axis.tick_labels.font.size = Pt(9)
    value_axis.tick_labels.number_format = value_axis_format
    if y_min is not None:
        value_axis.minimum_scale = y_min
    if y_max is not None:
        value_axis.maximum_scale = y_max
    for series, (_name, _values, color) in zip(chart.series, series_defs):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = color
        series.format.line.color.rgb = color


def add_tradeoff_scatter(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    points: list[dict[str, object]],
    *,
    x_min: float,
    x_max: float,
    y_min: float,
    y_max: float,
) -> None:
    plot_x, plot_y, plot_w, plot_h = add_chart_panel(
        slide,
        x,
        y,
        w,
        h,
        title,
        subtitle="横轴为 runtime 相对 baseline 变化，纵轴为 avg power 相对 baseline 变化",
    )

    left_margin = 0.70
    right_margin = 0.35
    top_margin = 0.20
    bottom_margin = 0.52
    left = plot_x + left_margin
    top = plot_y + top_margin
    width = plot_w - left_margin - right_margin
    height = plot_h - top_margin - bottom_margin

    chart_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    chart_bg.fill.solid()
    chart_bg.fill.fore_color.rgb = BG
    chart_bg.line.color.rgb = BORDER

    def project_x(value: float) -> float:
        return left + (value - x_min) / (x_max - x_min) * width

    def project_y(value: float) -> float:
        return top + height - (value - y_min) / (y_max - y_min) * height

    x_ticks = [-3.0, -1.5, 0.0, 1.5, 3.0]
    y_ticks = [-40.0, -30.0, -20.0, -10.0, 0.0]

    for tick in x_ticks:
        x_pos = project_x(tick)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x_pos), Inches(top), Inches(0.01), Inches(height))
        line.fill.solid()
        line.fill.fore_color.rgb = MUTED if abs(tick) < 1e-9 else BORDER
        line.line.fill.background()
        add_textbox(slide, x_pos - 0.18, top + height + 0.05, 0.36, 0.18, f"{tick:g}", size=8, color=MUTED, align=PP_ALIGN.CENTER)

    for tick in y_ticks:
        y_pos = project_y(tick)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(y_pos), Inches(width), Inches(0.01))
        line.fill.solid()
        line.fill.fore_color.rgb = MUTED if abs(tick) < 1e-9 else BORDER
        line.line.fill.background()
        add_textbox(slide, left - 0.52, y_pos - 0.09, 0.42, 0.18, f"{tick:g}", size=8, color=MUTED, align=PP_ALIGN.RIGHT)

    add_textbox(slide, left + width / 2 - 1.3, top + height + 0.28, 2.6, 0.18, "ΔRuntime vs Baseline (%)", size=9, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, left - 0.7, top - 0.02, 0.5, 0.55, "ΔPower\n(%)", size=9, color=MUTED, align=PP_ALIGN.CENTER)

    quadrant_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + w - 1.95), Inches(y + 0.15), Inches(1.65), Inches(0.75))
    quadrant_box.fill.solid()
    quadrant_box.fill.fore_color.rgb = SUCCESS_LIGHT
    quadrant_box.line.color.rgb = BORDER
    add_textbox(slide, x + w - 1.83, y + 0.23, 1.4, 0.14, "左下区域", size=9, bold=True, color=SUCCESS)
    add_textbox(slide, x + w - 1.83, y + 0.42, 1.4, 0.18, "更快且更省电", size=8, color=MUTED)

    legend_items = [("案例 A", ACCENT), ("案例 B", WARNING)]
    for idx, (label, color) in enumerate(legend_items):
        lx = x + 0.25 + idx * 1.35
        ly = y + 0.15
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(lx), Inches(ly), Inches(0.14), Inches(0.14))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.color.rgb = color
        add_textbox(slide, lx + 0.18, ly - 0.03, 0.9, 0.18, label, size=8, color=MUTED)

    for point in points:
        px = project_x(float(point["x"]))
        py = project_y(float(point["y"]))
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(px - 0.07), Inches(py - 0.07), Inches(0.14), Inches(0.14))
        circle.fill.solid()
        circle.fill.fore_color.rgb = point["color"]
        circle.line.color.rgb = point["color"]
        add_textbox(
            slide,
            px + float(point.get("dx", 0.08)),
            py + float(point.get("dy", -0.10)),
            0.95,
            0.18,
            str(point["label"]),
            size=8,
            color=TEXT,
        )


def axis_bounds(values: list[float], *, step: float = 5.0, padding: float = 5.0) -> tuple[float, float]:
    lower = min(values)
    upper = max(values)
    axis_min = min(0.0, math.floor((lower - padding) / step) * step)
    axis_max = max(0.0, math.ceil((upper + padding) / step) * step)
    return axis_min, axis_max


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_textbox(
        slide,
        1.0,
        2.0,
        14.0,
        1.4,
        "分布式训练 GPU 能效优化",
        size=30,
        bold=True,
        color=TITLE,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        1.0,
        3.0,
        14.0,
        1.0,
        "基于受控锁频、Zeus 统计与网络感知预测的实验研究",
        size=18,
        color=ACCENT,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        2.0,
        4.2,
        12.0,
        1.0,
        "Megatron-DeepSpeed baseline 对照 | V100 InfiniBand + RTX 4080 Ethernet",
        size=14,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    add_note_box(
        slide,
        3.2,
        5.5,
        9.6,
        1.6,
        "汇报重点",
        [
            "先给出 baseline 与 static 的受控对照证据，再给出 predictor 的跨拓扑结果",
            "所有 headline 均明确标注支撑来源与适用边界",
        ],
        fill_color=ACCENT_LIGHT,
    )
    add_footer(slide, "Academic-style internal report deck generated from local repo materials", 1)


def add_problem_slide(prs: Presentation, page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, "1. 研究问题与汇报主线", "先回答“为什么比较 baseline 与 static”，再回答“预测层解决什么问题”")
    add_bullets(
        slide,
        0.9,
        1.7,
        7.2,
        4.8,
        [
            "Megatron-DeepSpeed baseline 是默认训练路径，但默认 GPU 行为不一定处在能效最优点。",
            "本项目真正关注的问题不是“能不能降频”，而是“能否在不显著拉长训练时间的前提下降低功耗与能耗”。",
            "因此汇报必须先给出 baseline vs static 的受控对照证据，再讨论 predictor 如何帮助缩小 sweep 成本。",
        ],
        size=18,
    )
    add_note_box(
        slide,
        8.5,
        1.7,
        6.3,
        2.6,
        "本次 deck 的写法",
        [
            "不以“20%+”大字报开场",
            "先给实验口径，再给数据表，再给结论",
            "把 predictor 明确定位为 frequency-selection assistant",
        ],
        fill_color=SUCCESS_LIGHT,
        border_color=SUCCESS,
    )
    add_note_box(
        slide,
        8.5,
        4.7,
        6.3,
        1.8,
        "核心研究问题",
        [
            "Q1: 合适固定频点是否确实存在？",
            "Q2: 如何减少在新拓扑/新网络上的试错频点数？",
        ],
        fill_color=WARNING_LIGHT,
        border_color=WARNING,
    )
    add_footer(slide, "Narrative source: 汇报总结_20260415/01_实验口径与主线.md", page_no)


def add_protocol_slide(prs: Presentation, page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, "2. 公平对比口径：除了时钟策略，其他都不变")
    add_note_box(
        slide,
        0.9,
        1.55,
        6.9,
        2.1,
        "唯一允许变化的变量",
        [
            "`EXPERIMENT_MODE=baseline` 与 `EXPERIMENT_MODE=static`",
            "`STATIC_CLOCK_MHZ`",
        ],
        fill_color=ACCENT_LIGHT,
    )
    add_note_box(
        slide,
        8.2,
        1.55,
        6.9,
        2.1,
        "必须保持一致的条件",
        [
            "模型、tokenizer、数据集、TP/PP/DP、节点数、batch、训练步数、ZeRO、precision",
        ],
        fill_color=SUCCESS_LIGHT,
        border_color=SUCCESS,
    )
    build_table(
        slide,
        0.9,
        4.1,
        14.2,
        2.5,
        [
            ["模式", "时钟策略", "观测来源", "核心指标"],
            ["baseline", "默认 GPU 行为", "Zeus + run artifacts", "time / avg_power / energy"],
            ["static", "固定到目标 `STATIC_CLOCK_MHZ`", "Zeus + run artifacts", "time / avg_power / energy"],
        ],
        [2.2, 4.0, 4.0, 4.0],
    )
    add_bullets(
        slide,
        1.0,
        6.95,
        14.0,
        0.9,
        [
            "因此后文所有“功耗节省”都应理解为：在相同 workload 与拓扑下，相对 baseline 的受控变化。",
        ],
        size=14,
    )
    add_footer(slide, "Method rule source: docs/experiment-tracking.md and 汇报总结_20260415/05_证据清单.md", page_no)


def add_method_slide(prs: Presentation, page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, "3. 系统与方法：实验层、观测层、预测层")
    box_specs = [
        (
            0.9,
            "实验层",
            [
                "`scripts/run_experiment.sh` 统一 baseline/static 启动",
                "preflight + topology + hostfile snapshot 固化实验条件",
            ],
            ACCENT_LIGHT,
            ACCENT,
        ),
        (
            5.4,
            "观测层",
            [
                "Zeus 统一记录时间、平均功率、总能耗",
                "run.json / events.jsonl 作为实验工件",
            ],
            SUCCESS_LIGHT,
            SUCCESS,
        ),
        (
            9.9,
            "预测层",
            [
                "network benchmark 驱动连续 alpha 缩放",
                "cluster-capacity 与 per-node power scaling 处理拓扑迁移",
            ],
            WARNING_LIGHT,
            WARNING,
        ),
    ]
    for x, title, bullets, fill, border in box_specs:
        add_note_box(slide, x, 2.0, 4.0, 3.0, title, bullets, fill_color=fill, border_color=border)

    formula = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(5.5), Inches(13.6), Inches(1.4))
    formula.fill.solid()
    formula.fill.fore_color.rgb = RGBColor(255, 255, 255)
    formula.line.color.rgb = BORDER
    add_textbox(slide, 1.5, 5.75, 13.0, 0.3, "跨节点时间近似写作：T_target = T_base + α_tp·V_tp + α_pp·V_pp + α_dp·V_dp", size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, 1.5, 6.18, 13.0, 0.28, "关键点不是直接用 benchmark 时间替代预测，而是用 benchmark 推断当前拓扑的通信质量，再连续调整 cross-node penalty 参数。", size=11, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, "Method source: README.md, analysis/freq_model/*, 汇报总结_20260415/07_实现说明.md", page_no)


def add_evidence_slide(prs: Presentation, page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, "4. 实验环境与证据来源", "汇报中明确区分 artifact-backed 本地结果与 memory-bank preserved 历史摘要")
    build_table(
        slide,
        0.8,
        1.7,
        14.5,
        3.6,
        [
            ["部分", "平台 / 拓扑", "用途", "关键指标", "证据等级"],
            ["Baseline vs Static 案例 A", "V100, TP=1/PP=4/DP=4", "受控能效对照", "time / power / energy", "B"],
            ["Baseline vs Static 案例 B", "V100, TP=2/PP=2/DP=4", "受控能效对照", "time / power / energy", "B+A"],
            ["IB formal replay", "V100, 2x4 -> 2x8, IB", "predictor 准确性", "MAPE + per-point replay", "A"],
            ["Ethernet formal replay", "RTX 4080, 1x4 -> 2x4, Eth", "predictor 泛化", "MAPE + per-point replay", "A"],
        ],
        [3.0, 3.8, 2.4, 2.8, 1.5],
    )
    add_note_box(
        slide,
        0.9,
        5.8,
        6.9,
        1.65,
        "为什么要写证据等级",
        [
            "避免把历史摘要误说成“当前本地完整工件”",
            "避免把不同拓扑/不同网络的结果混写成同一组结论",
        ],
        fill_color=WARNING_LIGHT,
        border_color=WARNING,
    )
    add_note_box(
        slide,
        8.2,
        5.8,
        6.9,
        1.65,
        "本次 deck 的学术化改写原则",
        [
            "所有 headline 前面都有支撑页",
            "所有 predictor 结果都明确写明 topology 与 transport",
        ],
        fill_color=SUCCESS_LIGHT,
        border_color=SUCCESS,
    )
    add_footer(slide, "Evidence source: 汇报总结_20260415/05_证据清单.md, .context/paper/experimental_data.md", page_no)


def add_controlled_case_slide(prs: Presentation, page_no: int, case: ControlledCase) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, case.title, case.topology)

    baseline = case.baseline
    time_changes = [pct_change(p.time_s, baseline.time_s) for p in case.static_points]
    power_changes = [pct_change(p.avg_power_w, baseline.avg_power_w) for p in case.static_points]
    energy_changes = [pct_change(p.energy_j, baseline.energy_j) for p in case.static_points]
    y_min, y_max = axis_bounds(time_changes + power_changes + energy_changes)

    add_column_chart(
        slide,
        0.75,
        1.7,
        8.55,
        4.85,
        "相对 baseline 的变化 (%)",
        [point.label for point in case.static_points],
        [
            ("ΔTime", time_changes, ACCENT),
            ("ΔPower", power_changes, SUCCESS),
            ("ΔEnergy", energy_changes, WARNING),
        ],
        subtitle="负值更优；图中所有频点都与同一 baseline 对比",
        value_axis_format='0.0"%"',
        y_min=y_min,
        y_max=y_max,
        show_legend=True,
    )

    add_note_box(
        slide,
        9.6,
        1.7,
        5.1,
        1.6,
        "Baseline 锚点",
        [
            f"time = {baseline.time_s:.1f} s",
            f"avg power = {baseline.avg_power_w:.1f} W",
            f"energy = {fmt_energy_kj(baseline.energy_j)} kJ",
        ],
        title_size=13,
        bullet_size=11,
        fill_color=ACCENT_LIGHT,
        border_color=ACCENT,
    )

    delta_rows = [["频点", "ΔTime", "ΔPower", "ΔEnergy"]]
    for point in case.static_points:
        delta_rows.append(
            [
                point.label,
                fmt_pct(pct_change(point.time_s, baseline.time_s)),
                fmt_pct(pct_change(point.avg_power_w, baseline.avg_power_w)),
                fmt_pct(pct_change(point.energy_j, baseline.energy_j)),
            ]
        )

    build_table(
        slide,
        9.45,
        3.45,
        5.4,
        2.95,
        delta_rows,
        [1.9, 1.15, 1.15, 1.15],
        header_font_size=10,
        body_font_size=11,
    )

    best_energy_point = min(case.static_points, key=lambda point: point.energy_j)
    best_time_change = pct_change(best_energy_point.time_s, baseline.time_s)
    best_power_change = pct_change(best_energy_point.avg_power_w, baseline.avg_power_w)
    best_energy_change = pct_change(best_energy_point.energy_j, baseline.energy_j)
    observations = [
        f"runtime 变化范围：{fmt_pct(min(time_changes))} 到 {fmt_pct(max(time_changes))}",
        f"平均功率 / 总能耗：{fmt_pct(min(power_changes))}~{fmt_pct(max(power_changes))} / {fmt_pct(min(energy_changes))}~{fmt_pct(max(energy_changes))}",
    ]
    add_note_box(
        slide,
        0.9,
        6.72,
        7.0,
        1.38,
        "关键观察",
        observations,
        title_size=13,
        bullet_size=11,
        fill_color=SUCCESS_LIGHT,
        border_color=SUCCESS,
    )
    add_note_box(
        slide,
        8.2,
        6.72,
        6.9,
        1.38,
        "最佳折中点与支撑范围",
        [
            f"{best_energy_point.label}: {fmt_pct(best_time_change)} / {fmt_pct(best_power_change)} / {fmt_pct(best_energy_change)}",
            case.evidence_label,
        ],
        title_size=13,
        bullet_size=11,
        fill_color=ACCENT_LIGHT,
        border_color=ACCENT,
    )
    add_footer(slide, "Comparison source: 汇报总结_20260415/05_证据清单.md", page_no)


def add_controlled_summary_slide(prs: Presentation, page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, "7. 受控对照实验的可支撑结论")
    add_metric_card(slide, 0.9, 1.7, 3.2, 1.9, "-2.4% ~ +2.2%", "runtime 相对 baseline", "时间大多保持在 baseline 附近", value_color=ACCENT)
    add_metric_card(slide, 4.4, 1.7, 3.2, 1.9, "-24.1% ~ -36.3%", "avg power 相对 baseline", "大部分点显著落在省电区间", value_color=SUCCESS)
    add_metric_card(slide, 7.9, 1.7, 3.2, 1.9, "-25.8% ~ -34.9%", "energy 相对 baseline", "总能耗与功率结论方向一致", value_color=SUCCESS)
    add_metric_card(slide, 11.4, 1.7, 3.2, 1.9, "7 points", "当前主要支撑样例", "均为受控 baseline/static 对照", value_color=ACCENT)

    scatter_points = [
        {"x": -1.3, "y": -25.3, "label": "A-1252", "color": ACCENT, "dx": 0.10, "dy": -0.18},
        {"x": -2.2, "y": -24.9, "label": "A-1260", "color": ACCENT, "dx": 0.08, "dy": 0.02},
        {"x": -2.3, "y": -24.1, "label": "A-1267", "color": ACCENT, "dx": 0.08, "dy": -0.28},
        {"x": 2.2, "y": -36.3, "label": "B-1072", "color": WARNING, "dx": 0.08, "dy": -0.18},
        {"x": 1.5, "y": -35.9, "label": "B-1080", "color": WARNING, "dx": 0.08, "dy": 0.00},
        {"x": 0.9, "y": -35.3, "label": "B-1087", "color": WARNING, "dx": 0.08, "dy": -0.22},
        {"x": -2.4, "y": -32.8, "label": "B-1125", "color": WARNING, "dx": 0.08, "dy": 0.04},
    ]
    add_tradeoff_scatter(
        slide,
        0.8,
        3.95,
        9.45,
        3.95,
        "时间-功率 trade-off：所有 static 点相对 baseline 的位置",
        scatter_points,
        x_min=-3.0,
        x_max=3.0,
        y_min=-40.0,
        y_max=0.0,
    )

    add_note_box(
        slide,
        10.55,
        4.05,
        4.45,
        1.65,
        "因此可以稳妥写出的 headline",
        [
            "在相同 workload 与拓扑下，合适固定频点可带来约 20% 量级的平均功率下降。",
            "在已完成的 V100 双机案例中，部分拓扑可达到 25% 到 35%+ 的平均功率下降，同时 runtime 基本不变。",
        ],
        title_size=13,
        bullet_size=11,
        fill_color=SUCCESS_LIGHT,
        border_color=SUCCESS,
    )
    add_note_box(
        slide,
        10.55,
        6.0,
        4.45,
        1.65,
        "当前不能直接说死的内容",
        [
            "不能说 predictor 本身直接带来了 20% 节能。",
            "不能说所有拓扑都会落在同一个最佳频点。",
            "不能把所有历史 baseline 工件都描述成已本地完整保留。",
        ],
        title_size=13,
        bullet_size=10,
        fill_color=WARNING_LIGHT,
        border_color=WARNING,
    )
    add_footer(slide, "Headline source: 汇报总结_20260415/01_实验口径与主线.md and 05_证据清单.md", page_no)


def add_predictor_slide(prs: Presentation, page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, "8. 为什么还需要 predictor：减少 sweep 成本，而不是替代实验")
    add_note_box(
        slide,
        0.9,
        1.7,
        6.7,
        2.5,
        "问题",
        [
            "真实 baseline/static 对照可以证明节能，但每到新拓扑、新网络都完整 sweep 成本很高。",
            "跨节点训练的主要困难在于通信罚时并不是固定常数，网络和拓扑变化都会影响它。",
        ],
        fill_color=ACCENT_LIGHT,
        border_color=ACCENT,
    )
    add_note_box(
        slide,
        8.0,
        1.7,
        7.1,
        2.5,
        "当前 predictor 的核心改动",
        [
            "用 all-reduce benchmark 估计当前拓扑的通信质量，并连续缩放 alpha/beta 参数。",
            "增加 cluster-capacity scaling 与 per-node GPU-count power scaling，避免直接把 source anchor 原样套到 target。",
        ],
        fill_color=SUCCESS_LIGHT,
        border_color=SUCCESS,
    )
    formula = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.7), Inches(14.1), Inches(1.4))
    formula.fill.solid()
    formula.fill.fore_color.rgb = RGBColor(255, 255, 255)
    formula.line.color.rgb = BORDER
    add_textbox(slide, 1.2, 4.98, 13.7, 0.3, "predictor 的定位：frequency-selection assistant", size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, 1.2, 5.4, 13.7, 0.3, "真实证据仍来自 baseline/static + Zeus；predictor 的价值是减少值得验证的频点数量，并提高跨拓扑迁移效率。", size=12, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, "Method scope: README.md, 汇报总结_20260415/01_实验口径与主线.md, analysis/freq_model/*", page_no)


def add_ib_slide(prs: Presentation, page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, "9. IB formal replay：当前 predictor 的主要精度证据", IB_METRICS["pair"])
    add_metric_card(slide, 0.9, 1.6, 3.6, 1.7, f"{IB_METRICS['time_mape']:.2f}%", "Time MAPE", "fresh transport-consistent replay", value_color=ACCENT)
    add_metric_card(slide, 4.8, 1.6, 3.6, 1.7, f"{IB_METRICS['power_mape']:.2f}%", "Power MAPE", "after per-node power scaling", value_color=SUCCESS)
    add_metric_card(slide, 8.7, 1.6, 3.6, 1.7, f"{IB_METRICS['energy_mape']:.2f}%", "Energy MAPE", "artifact-backed local replay", value_color=SUCCESS)
    add_metric_card(slide, 12.9, 1.6, 2.2, 1.7, IB_METRICS["alpha_dp"], "alpha_dp", "IB fast-network scale", value_color=ACCENT)
    ib_labels = [row[0] for row in IB_METRICS["points"]]
    add_column_chart(
        slide,
        0.8,
        3.65,
        7.0,
        3.15,
        "时间：观测值 vs 预测值",
        ib_labels,
        [
            ("Observed", [float(row[1]) for row in IB_METRICS["points"]], ACCENT),
            ("Predicted", [float(row[2]) for row in IB_METRICS["points"]], WARNING),
        ],
        subtitle="当前主要剩余误差集中在 runtime",
        value_axis_format="0.0",
        show_legend=True,
    )
    add_column_chart(
        slide,
        8.1,
        3.65,
        7.0,
        3.15,
        "平均功率：观测值 vs 预测值",
        ib_labels,
        [
            ("Observed", [float(row[4]) for row in IB_METRICS["points"]], SUCCESS),
            ("Predicted", [float(row[5]) for row in IB_METRICS["points"]], ACCENT),
        ],
        subtitle="power path 已明显收敛",
        value_axis_format="0",
        show_legend=True,
    )
    add_note_box(
        slide,
        0.9,
        6.95,
        14.2,
        1.0,
        "解释",
        [
            "旧的 `98.5%` 是历史 motivating failure；当前正式 paired replay 已收敛到 `11.48% / 3.28% / 7.86%`。剩余主要误差在 runtime，而不再是 power/energy。",
        ],
        title_size=13,
        bullet_size=11,
        fill_color=ACCENT_LIGHT,
        border_color=ACCENT,
    )
    add_footer(slide, "Artifact source: .context/paper/experimental_data.md and IB transfer_prediction_report.md", page_no)


def add_eth_slide(prs: Presentation, page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, "10. Ethernet formal replay：预测层在慢网络场景仍可工作", ETH_METRICS["pair"])
    add_metric_card(slide, 0.9, 1.6, 3.4, 1.7, f"{ETH_METRICS['time_mape']:.2f}%", "Time MAPE", "formal `1x4 -> 2x4` replay", value_color=ACCENT)
    add_metric_card(slide, 4.6, 1.6, 3.4, 1.7, f"{ETH_METRICS['power_mape']:.2f}%", "Power MAPE", "power 侧仍明显弱于 IB", value_color=WARNING)
    add_metric_card(slide, 8.3, 1.6, 3.4, 1.7, f"{ETH_METRICS['energy_mape']:.2f}%", "Energy MAPE", "energy 已接近可比较边界", value_color=SUCCESS)
    add_metric_card(slide, 12.0, 1.6, 3.1, 1.7, ETH_METRICS["alpha_dp"], "alpha_dp", "slow-network scale", value_color=WARNING)
    eth_labels = [row[0] for row in ETH_METRICS["points"]]
    add_column_chart(
        slide,
        0.8,
        3.65,
        7.0,
        3.15,
        "时间：观测值 vs 预测值",
        eth_labels,
        [
            ("Observed", [float(row[1]) for row in ETH_METRICS["points"]], ACCENT),
            ("Predicted", [float(row[2]) for row in ETH_METRICS["points"]], WARNING),
        ],
        subtitle="time path 已进入可用范围",
        value_axis_format="0.0",
        show_legend=True,
    )
    add_column_chart(
        slide,
        8.1,
        3.65,
        7.0,
        3.15,
        "逐频点 APE (%)",
        eth_labels,
        [
            ("Time APE", [float(row[3].rstrip('%')) for row in ETH_METRICS["points"]], ACCENT),
            ("Power APE", [float(row[4].rstrip('%')) for row in ETH_METRICS["points"]], WARNING),
        ],
        subtitle="power 侧仍明显弱于 IB",
        value_axis_format='0.0"%"',
        y_min=0.0,
        y_max=20.0,
        show_legend=True,
    )
    add_note_box(
        slide,
        0.9,
        6.95,
        14.2,
        1.0,
        "边界说明",
        [
            "这组 Ethernet 结果的 topology 是 `TP=1, PP=2`，与 IB 主结果的 `TP=4, PP=1` 不同，因此二者的 MAPE 不能直接当作 apples-to-apples 排名。它更适合支撑“network-aware predictor 可迁移到慢网络场景”这一结论。",
        ],
        title_size=13,
        bullet_size=10,
        fill_color=WARNING_LIGHT,
        border_color=WARNING,
    )
    add_footer(slide, "Artifact source: Ethernet transfer_prediction_report.md and memory-bank/activeContext.md", page_no)


def add_conclusion_slide(prs: Presentation, page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, "11. 结论与边界：这份汇报真正证明了什么")
    add_note_box(
        slide,
        0.9,
        1.65,
        6.9,
        3.5,
        "可以稳妥主张的内容",
        [
            "Megatron-DeepSpeed baseline 不是天然的能效最优点。",
            "在受控 baseline/static 对照下，合适固定频点可以在几乎不改变训练时间的前提下降低平均功率与总能耗。",
            "predictor 在 IB 与 Ethernet 两类网络下都能提供可用的跨拓扑时间预测，并帮助缩小候选频点范围。",
        ],
        fill_color=SUCCESS_LIGHT,
        border_color=SUCCESS,
    )
    add_note_box(
        slide,
        8.1,
        1.65,
        7.0,
        3.5,
        "当前仍需避免的过强表述",
        [
            "不要把所有“20%+ 节能”都归功于 predictor。",
            "不要把不同 topology / transport 的 MAPE 直接排成单一排行榜。",
            "不要把历史 preserved summary 说成全部是当前本地完整原始工件。",
        ],
        fill_color=WARNING_LIGHT,
        border_color=WARNING,
    )
    add_note_box(
        slide,
        0.9,
        5.55,
        14.2,
        1.45,
        "下一步",
        [
            "继续补齐更多 transport-consistent 复现实验，并把当前图表化证据进一步下沉到论文正式图与答辩材料中。",
        ],
        fill_color=ACCENT_LIGHT,
        border_color=ACCENT,
    )
    add_footer(slide, "Conclusion source: 汇报总结_20260415/01_实验口径与主线.md, 05_证据清单.md, .context/paper/experimental_data.md", page_no)


def add_thanks_slide(prs: Presentation, page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_textbox(slide, 1.0, 2.3, 14.0, 0.9, "谢谢", size=34, bold=True, color=TITLE, align=PP_ALIGN.CENTER)
    add_textbox(slide, 1.0, 3.25, 14.0, 0.5, "Q&A", size=18, color=ACCENT, align=PP_ALIGN.CENTER)
    add_note_box(
        slide,
        3.2,
        4.7,
        9.6,
        1.7,
        "附记",
        [
            "当前 deck 已经把 baseline/static 受控证据与 predictor 结果分层展示。若继续打磨，优先补图，不优先堆更多 headline。",
        ],
        fill_color=ACCENT_LIGHT,
        border_color=ACCENT,
    )
    add_footer(slide, "Generated by 汇报总结_20260415/generate_ppt.py", page_no)


def create_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    add_title_slide(prs)
    add_problem_slide(prs, 2)
    add_protocol_slide(prs, 3)
    add_method_slide(prs, 4)
    add_evidence_slide(prs, 5)
    add_controlled_case_slide(prs, 6, CASE_A)
    add_controlled_case_slide(prs, 7, CASE_B)
    add_controlled_summary_slide(prs, 8)
    add_predictor_slide(prs, 9)
    add_ib_slide(prs, 10)
    add_eth_slide(prs, 11)
    add_conclusion_slide(prs, 12)
    add_thanks_slide(prs, 13)
    return prs


if __name__ == "__main__":
    presentation = create_presentation()
    output_path = "GPU能效优化_20%+功耗节省方案.pptx"
    presentation.save(output_path)
    print(f"Generated: {output_path}")
    print(f"Slides: {len(presentation.slides)}")
